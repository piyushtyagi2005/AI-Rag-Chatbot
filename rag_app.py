"""
AI Chatbot — Streamlit UI with RAG + Conversation Memory
Gen Z edition: neon gradients, glassmorphism, vibrant palette.

Run:
    pip install streamlit requests pypdf python-docx openpyxl python-pptx python-dotenv
    streamlit run stackroom_app.py
"""

import os, re, math, io, csv
from collections import defaultdict

import streamlit as st
import requests

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# ── Config ────────────────────────────────────────────────────────────────────
GEMINI_MODEL  = os.environ.get("GEMINI_MODEL", "gemini-2.5-flash")
GEMINI_URL    = (
    f"https://generativelanguage.googleapis.com/v1beta/models/"
    f"{GEMINI_MODEL}:generateContent"
)
CHUNK_SIZE    = 400
CHUNK_OVERLAP = 50
TOP_K         = 4

# ═══════════════════════════════════════════════════════════════════════════════
#  TF-IDF RAG engine
# ═══════════════════════════════════════════════════════════════════════════════

def _tokenize(text):
    return re.findall(r"[a-z0-9]+", text.lower())

def _chunk_text(text, size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    words = text.split()
    step  = max(1, size - overlap)
    return [" ".join(words[i:i+size]) for i in range(0, len(words), step) if words[i:i+size]]

def _build_tfidf(chunks):
    n = len(chunks)
    tf_list, df = [], defaultdict(int)
    for chunk in chunks:
        tokens = _tokenize(chunk)
        total  = max(len(tokens), 1)
        tf     = defaultdict(float)
        for t in tokens:
            tf[t] += 1.0 / total
        tf_list.append(dict(tf))
        for term in set(tokens):
            df[term] += 1
    idf   = {t: math.log((n+1)/(c+1))+1.0 for t, c in df.items()}
    tfidf = [{t: w * idf.get(t, 1.0) for t, w in tf.items()} for tf in tf_list]
    return tfidf, idf

def _cosine(a, b):
    if not a or not b: return 0.0
    dot    = sum(a.get(t, 0.0) * w for t, w in b.items())
    norm_a = math.sqrt(sum(v*v for v in a.values()))
    norm_b = math.sqrt(sum(v*v for v in b.values()))
    denom  = norm_a * norm_b
    return dot / denom if denom else 0.0

def _query_rag(query):
    """
    Multi-document aware retrieval.

    Strategy:
    1. Score every chunk against the query using cosine similarity.
    2. From EACH document pick its single best-scoring chunk  →  guarantees
       every uploaded doc contributes at least one result.
    3. Fill any remaining slots (up to top_k) with the next highest-scoring
       chunks from the global ranked list (may come from any doc).
    4. De-duplicate by chunk index and return sorted by score.
    """
    chunks = st.session_state.rag_chunks
    tfidf  = st.session_state.rag_tfidf
    idf    = st.session_state.rag_idf
    top_k  = st.session_state.get("top_k", TOP_K)

    if not chunks:
        return []

    # ── Build query vector ────────────────────────────────────────────────
    tokens   = _tokenize(query)
    total    = max(len(tokens), 1)
    query_tf = defaultdict(float)
    for t in tokens:
        query_tf[t] += 1.0 / total
    qvec = {t: w * idf.get(t, 1.0) for t, w in query_tf.items()}

    # ── Score all chunks ──────────────────────────────────────────────────
    scored = sorted(
        [(_cosine(qvec, v), i) for i, v in enumerate(tfidf)],
        reverse=True
    )

    # ── Step 1: best chunk per document (guaranteed representation) ───────
    best_per_doc: dict[str, tuple] = {}   # doc_name → (score, chunk_idx)
    for score, idx in scored:
        doc = chunks[idx]["doc"]
        if doc not in best_per_doc:
            best_per_doc[doc] = (score, idx)

    selected_indices = {idx for _, idx in best_per_doc.values()}

    # ── Step 2: fill remaining slots from global top list ─────────────────
    for score, idx in scored:
        if len(selected_indices) >= top_k:
            break
        if idx not in selected_indices and score > 0:
            selected_indices.add(idx)

    # ── Step 3: build result list sorted by score ─────────────────────────
    results = []
    for score, idx in scored:
        if idx in selected_indices and score > 0:
            results.append({"score": score, **chunks[idx]})

    return results

def _rebuild_index():
    texts = [c["text"] for c in st.session_state.rag_chunks]
    if not texts:
        st.session_state.rag_tfidf = []
        st.session_state.rag_idf   = {}
        return
    tfidf, idf = _build_tfidf(texts)
    st.session_state.rag_tfidf = tfidf
    st.session_state.rag_idf   = idf

def ingest_text(filename, text):
    st.session_state.rag_chunks = [c for c in st.session_state.rag_chunks if c["doc"] != filename]
    new_chunks = _chunk_text(text)
    for ch in new_chunks:
        st.session_state.rag_chunks.append({"doc": filename, "text": ch})
    if filename not in st.session_state.doc_list:
        st.session_state.doc_list.append(filename)
    _rebuild_index()
    return len(new_chunks)

# ═══════════════════════════════════════════════════════════════════════════════
#  File extraction
# ═══════════════════════════════════════════════════════════════════════════════

def extract_text(file):
    _, ext = os.path.splitext(file.name.lower())
    raw = file.read()
    if ext in [".txt",".md",".json",".xml",".html",".htm",".py",".js",".css"]:
        return raw.decode("utf-8", errors="ignore")
    elif ext == ".pdf":
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(raw))
        return "\n".join(p for page in reader.pages if (p := page.extract_text() or ""))
    elif ext in [".docx",".doc"]:
        import docx
        doc = docx.Document(io.BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs)
    elif ext in [".xlsx",".xls"]:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            parts.append(f"--- Sheet: {sheet} ---")
            for row in wb[sheet].iter_rows(values_only=True):
                r = ", ".join(str(c) if c is not None else "" for c in row)
                if r.strip(): parts.append(r)
        return "\n".join(parts)
    elif ext == ".csv":
        stream = io.StringIO(raw.decode("utf-8", errors="ignore"), newline="")
        return "\n".join(", ".join(row) for row in csv.reader(stream) if any(row))
    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)
    else:
        return raw.decode("utf-8", errors="ignore")

# ═══════════════════════════════════════════════════════════════════════════════
#  Gemini
# ═══════════════════════════════════════════════════════════════════════════════

def build_gemini_contents(messages):
    contents = []
    for msg in messages:
        grole = "model" if msg["role"] == "assistant" else "user"
        text  = msg["content"] if isinstance(msg["content"], str) else str(msg["content"])
        if not text.strip(): continue
        if contents and contents[-1]["role"] == grole:
            contents[-1]["parts"][0]["text"] += "\n" + text
        else:
            contents.append({"role": grole, "parts": [{"text": text}]})
    while contents and contents[0]["role"] != "user":
        contents.pop(0)
    return contents

def call_gemini(messages, system, api_key):
    contents = build_gemini_contents(messages)
    if not contents: return "No valid messages."
    payload = {"contents": contents}
    if system:
        payload["system_instruction"] = {"parts": [{"text": system}]}
    try:
        resp = requests.post(
            GEMINI_URL,
            headers={"x-goog-api-key": api_key, "Content-Type": "application/json"},
            json=payload, timeout=60,
        )
        data = resp.json()
        if not resp.ok:
            return f"❌ {data.get('error',{}).get('message', resp.status_code)}"
        parts = data.get("candidates",[{}])[0].get("content",{}).get("parts",[])
        return "".join(p.get("text","") for p in parts) or "(empty response)"
    except Exception as e:
        return f"❌ {e}"

# ═══════════════════════════════════════════════════════════════════════════════
#  Session state
# ═══════════════════════════════════════════════════════════════════════════════

def init_state():
    for k, v in {
        "rag_chunks":   [],
        "rag_tfidf":    [],
        "rag_idf":      {},
        "doc_list":     [],
        "doc_meta":     {},
        "chat_history": [],
        "api_key":      os.environ.get("GEMINI_API_KEY",""),
        "show_sources": True,
        "show_chunks":  False,
        "top_k":        TOP_K,
    }.items():
        if k not in st.session_state:
            st.session_state[k] = v

# ═══════════════════════════════════════════════════════════════════════════════
#  Page config
# ═══════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="AI Chatbot",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Gen Z CSS ─────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap');

/* ── Root / body ── */
html, body, [data-testid="stAppViewContainer"] {
    background: #0a0a0f !important;
    font-family: 'Space Grotesk', sans-serif;
}

/* ── Sidebar ── */
[data-testid="stSidebar"] {
    background: linear-gradient(180deg, #0d0d1a 0%, #110d1f 100%) !important;
    border-right: 1px solid rgba(139,92,246,0.25) !important;
}
[data-testid="stSidebar"] > div:first-child { padding: 1rem 1rem 2rem; }

/* ── Main content bg ── */
[data-testid="stMainBlockContainer"] {
    background: transparent;
}
section[data-testid="stMain"] {
    background: radial-gradient(ellipse at 20% 0%, rgba(139,92,246,0.12) 0%, transparent 60%),
                radial-gradient(ellipse at 80% 100%, rgba(236,72,153,0.08) 0%, transparent 60%),
                #0a0a0f;
}

/* ── Headings ── */
h1, h2, h3 { font-family: 'Space Grotesk', sans-serif !important; }

/* ── Masthead ── */
.masthead-wrap {
    background: linear-gradient(135deg, rgba(139,92,246,0.15), rgba(236,72,153,0.1), rgba(6,182,212,0.08));
    border: 1px solid rgba(139,92,246,0.3);
    border-radius: 20px;
    padding: 22px 28px 18px;
    margin-bottom: 6px;
    position: relative;
    overflow: hidden;
}
.masthead-wrap::before {
    content: '';
    position: absolute;
    top: -50%; left: -50%;
    width: 200%; height: 200%;
    background: conic-gradient(from 0deg, transparent 0%, rgba(139,92,246,0.05) 25%, transparent 50%);
    animation: spin 8s linear infinite;
}
@keyframes spin { to { transform: rotate(360deg); } }
.masthead-title {
    font-size: 36px;
    font-weight: 700;
    background: linear-gradient(135deg, #a78bfa, #f472b6, #22d3ee);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    margin: 0;
    letter-spacing: -0.5px;
    position: relative;
}
.masthead-sub {
    font-family: 'JetBrains Mono', monospace;
    font-size: 11px;
    color: rgba(167,139,250,0.7);
    letter-spacing: 3px;
    text-transform: uppercase;
    margin-top: 4px;
    position: relative;
}

/* ── Neon divider ── */
.neon-divider {
    height: 1px;
    background: linear-gradient(90deg, transparent, #8b5cf6, #ec4899, #06b6d4, transparent);
    margin: 16px 0;
    border: none;
    opacity: 0.6;
}

/* ── Status badge ── */
.status-badge {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    padding: 5px 14px;
    border-radius: 999px;
    font-family: 'JetBrains Mono', monospace;
    font-size: 11px;
    font-weight: 500;
    letter-spacing: 1px;
}
.status-badge.rag-on {
    background: rgba(34,211,238,0.12);
    border: 1px solid rgba(34,211,238,0.4);
    color: #22d3ee;
}
.status-badge.rag-off {
    background: rgba(139,92,246,0.1);
    border: 1px solid rgba(139,92,246,0.3);
    color: rgba(167,139,250,0.7);
}
.status-dot {
    width: 6px; height: 6px;
    border-radius: 50%;
    animation: blink 1.5s ease-in-out infinite;
}
.rag-on .status-dot { background: #22d3ee; box-shadow: 0 0 6px #22d3ee; }
.rag-off .status-dot { background: #8b5cf6; }
@keyframes blink { 0%,100%{opacity:1;} 50%{opacity:0.3;} }

/* ── Chat bubbles ── */
.user-bubble {
    background: linear-gradient(135deg, #7c3aed, #a855f7);
    color: #fff;
    padding: 12px 18px;
    border-radius: 18px 18px 4px 18px;
    margin: 8px 0 8px 18%;
    font-size: 15px;
    line-height: 1.6;
    box-shadow: 0 4px 20px rgba(139,92,246,0.35);
    position: relative;
}
.user-bubble::before {
    content: '👤';
    position: absolute;
    right: -30px; top: 8px;
    font-size: 16px;
}

.bot-bubble {
    background: linear-gradient(135deg, rgba(15,15,30,0.95), rgba(20,10,35,0.95));
    border: 1px solid rgba(139,92,246,0.3);
    border-top: 2px solid;
    border-image: linear-gradient(90deg, #8b5cf6, #ec4899) 1;
    color: #e2e8f0;
    padding: 14px 18px;
    border-radius: 4px 18px 18px 18px;
    margin: 8px 18% 8px 0;
    font-size: 15px;
    line-height: 1.7;
    box-shadow: 0 4px 20px rgba(0,0,0,0.4), inset 0 1px 0 rgba(139,92,246,0.1);
    position: relative;
}
.bot-bubble::before {
    content: '🤖';
    position: absolute;
    left: -30px; top: 8px;
    font-size: 16px;
}
.bot-bubble.error {
    border-color: rgba(239,68,68,0.5);
    color: #fca5a5;
    font-family: 'JetBrains Mono', monospace;
    font-size: 13px;
}

/* ── Source chips ── */
.source-chips {
    display: flex;
    gap: 8px;
    flex-wrap: wrap;
    margin: 6px 18% 0 0;
    padding-left: 4px;
}
.chip {
    background: rgba(6,182,212,0.1);
    border: 1px solid rgba(6,182,212,0.35);
    color: #22d3ee;
    padding: 4px 12px;
    font-size: 11px;
    font-family: 'JetBrains Mono', monospace;
    border-radius: 999px;
    letter-spacing: 0.5px;
}

/* ── Doc cards ── */
.doc-card {
    background: rgba(139,92,246,0.08);
    border: 1px solid rgba(139,92,246,0.25);
    border-left: 3px solid #8b5cf6;
    padding: 10px 12px;
    border-radius: 10px;
    margin-bottom: 8px;
    color: #e2e8f0;
    font-size: 13px;
    transition: border-color 0.2s;
}
.doc-card:hover { border-color: rgba(236,72,153,0.5); }
.doc-card small {
    font-family: 'JetBrains Mono', monospace;
    font-size: 10px;
    color: rgba(167,139,250,0.6);
}

/* ── Memory badge ── */
.mem-badge {
    background: rgba(236,72,153,0.08);
    border: 1px solid rgba(236,72,153,0.25);
    border-radius: 999px;
    padding: 5px 14px;
    font-family: 'JetBrains Mono', monospace;
    font-size: 11px;
    color: rgba(244,114,182,0.8);
    display: inline-flex;
    align-items: center;
    gap: 7px;
}
.mem-dot {
    width: 7px; height: 7px;
    border-radius: 50%;
    background: #f472b6;
    box-shadow: 0 0 6px #f472b6;
    animation: blink 2s ease-in-out infinite;
}

/* ── Empty state ── */
.empty-state {
    text-align: center;
    padding: 70px 20px;
    color: rgba(139,92,246,0.5);
}
.empty-state .big-emoji { font-size: 52px; margin-bottom: 16px; }
.empty-state h3 {
    font-size: 22px;
    font-weight: 600;
    background: linear-gradient(135deg, #a78bfa, #f472b6);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    margin: 0 0 10px;
}
.empty-state p {
    font-size: 14px;
    color: rgba(226,232,240,0.4);
    max-width: 380px;
    margin: auto;
    line-height: 1.7;
}

/* ── Sidebar labels ── */
.sidebar-label {
    font-family: 'JetBrains Mono', monospace;
    font-size: 10px;
    letter-spacing: 2px;
    text-transform: uppercase;
    color: rgba(167,139,250,0.6);
    margin-bottom: 8px;
}

/* ── Empty catalog ── */
.empty-catalog {
    font-family: 'JetBrains Mono', monospace;
    font-size: 11px;
    color: rgba(139,92,246,0.35);
    border: 1px dashed rgba(139,92,246,0.2);
    padding: 14px;
    text-align: center;
    border-radius: 10px;
}

/* ── Streamlit overrides ── */
[data-testid="stChatInput"] > div {
    background: rgba(139,92,246,0.08) !important;
    border: 1px solid rgba(139,92,246,0.3) !important;
    border-radius: 14px !important;
}
[data-testid="stChatInput"] textarea {
    color: #e2e8f0 !important;
    font-family: 'Space Grotesk', sans-serif !important;
    font-size: 15px !important;
    background: transparent !important;
}
[data-testid="stChatInput"] textarea::placeholder { color: rgba(167,139,250,0.4) !important; }

/* Buttons */
.stButton > button {
    background: linear-gradient(135deg, rgba(139,92,246,0.2), rgba(236,72,153,0.15)) !important;
    border: 1px solid rgba(139,92,246,0.4) !important;
    color: #e2e8f0 !important;
    border-radius: 10px !important;
    font-family: 'Space Grotesk', sans-serif !important;
    font-weight: 500 !important;
    transition: all 0.2s !important;
}
.stButton > button:hover {
    background: linear-gradient(135deg, rgba(139,92,246,0.4), rgba(236,72,153,0.3)) !important;
    border-color: rgba(167,139,250,0.7) !important;
    box-shadow: 0 0 15px rgba(139,92,246,0.3) !important;
    transform: translateY(-1px) !important;
}

/* Inputs */
.stTextInput > div > input, .stTextArea > div > textarea {
    background: rgba(139,92,246,0.06) !important;
    border: 1px solid rgba(139,92,246,0.25) !important;
    border-radius: 10px !important;
    color: #e2e8f0 !important;
    font-family: 'Space Grotesk', sans-serif !important;
}
.stTextInput > div > input:focus, .stTextArea > div > textarea:focus {
    border-color: rgba(167,139,250,0.6) !important;
    box-shadow: 0 0 0 3px rgba(139,92,246,0.15) !important;
}

/* File uploader */
[data-testid="stFileUploader"] {
    background: rgba(139,92,246,0.05) !important;
    border: 1px dashed rgba(139,92,246,0.3) !important;
    border-radius: 12px !important;
}

/* Radio */
.stRadio > div { gap: 8px; }
.stRadio label {
    color: rgba(226,232,240,0.7) !important;
    font-family: 'Space Grotesk', sans-serif !important;
}

/* Select / expander */
.streamlit-expanderHeader {
    background: rgba(139,92,246,0.08) !important;
    border-radius: 10px !important;
    color: rgba(167,139,250,0.8) !important;
    font-family: 'JetBrains Mono', monospace !important;
    font-size: 12px !important;
}

/* Scrollbar */
::-webkit-scrollbar { width: 6px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: rgba(139,92,246,0.3); border-radius: 3px; }

/* Slider */
.stSlider > div > div { background: rgba(139,92,246,0.2) !important; }

/* Toggle */
.stToggle > label { color: rgba(226,232,240,0.7) !important; }

/* Success / error messages */
.stSuccess { background: rgba(34,197,94,0.1) !important; border-color: rgba(34,197,94,0.3) !important; color: #86efac !important; border-radius: 10px !important; }
.stError   { background: rgba(239,68,68,0.1) !important; border-color: rgba(239,68,68,0.3) !important; color: #fca5a5 !important; border-radius: 10px !important; }
.stWarning { background: rgba(251,191,36,0.1) !important; border-color: rgba(251,191,36,0.3) !important; color: #fde68a !important; border-radius: 10px !important; }

/* Divider */
hr { border-color: rgba(139,92,246,0.15) !important; }
</style>
""", unsafe_allow_html=True)

init_state()

# ═══════════════════════════════════════════════════════════════════════════════
#  Sidebar
# ═══════════════════════════════════════════════════════════════════════════════

with st.sidebar:
    # Masthead
    st.markdown("""
    <div style="margin-bottom:4px">
      <div style="font-size:28px;font-weight:700;background:linear-gradient(135deg,#a78bfa,#f472b6,#22d3ee);
           -webkit-background-clip:text;-webkit-text-fill-color:transparent;background-clip:text;
           letter-spacing:-0.5px">🤖 AI Chatbot</div>
      <div style="font-family:'JetBrains Mono',monospace;font-size:10px;letter-spacing:3px;
           color:rgba(167,139,250,0.6);text-transform:uppercase;margin-top:3px">
        RAG · Memory · Gemini
      </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown('<div class="neon-divider"></div>', unsafe_allow_html=True)

    # ── API Key ────────────────────────────────────────────────────────────
    with st.expander("🔑 API Key", expanded=not st.session_state.api_key):
        key_in = st.text_input("Gemini API Key", value=st.session_state.api_key,
                               type="password", label_visibility="collapsed",
                               placeholder="AIza…")
        if key_in != st.session_state.api_key:
            st.session_state.api_key = key_in
        if st.session_state.api_key:
            st.success("Key loaded ✓")
        else:
            st.warning("Get a free key → [aistudio.google.com](https://aistudio.google.com)")

    st.markdown('<div class="neon-divider"></div>', unsafe_allow_html=True)

    # ── Document catalog ───────────────────────────────────────────────────
    n_docs   = len(st.session_state.doc_list)
    n_chunks = len(st.session_state.rag_chunks)
    st.markdown(
        f'<div class="sidebar-label">📂 Knowledge Base &nbsp;'
        f'<span style="color:rgba(34,211,238,0.6)">{n_docs} docs · {n_chunks} chunks</span></div>',
        unsafe_allow_html=True
    )

    if st.session_state.doc_list:
        for fname in list(st.session_state.doc_list):
            meta   = st.session_state.doc_meta.get(fname, {})
            words  = meta.get("word_count", 0)
            chunks = meta.get("chunks", 0)
            c1, c2 = st.columns([5, 1])
            with c1:
                st.markdown(
                    f'<div class="doc-card">✦ {fname}<br>'
                    f'<small>{words:,} words · {chunks} chunks</small></div>',
                    unsafe_allow_html=True
                )
            with c2:
                st.markdown("<div style='padding-top:6px'>", unsafe_allow_html=True)
                if st.button("✕", key=f"del_{fname}", help=f"Remove {fname}"):
                    st.session_state.rag_chunks = [c for c in st.session_state.rag_chunks if c["doc"] != fname]
                    st.session_state.doc_list.remove(fname)
                    st.session_state.doc_meta.pop(fname, None)
                    _rebuild_index()
                    st.rerun()
                st.markdown("</div>", unsafe_allow_html=True)
    else:
        st.markdown('<div class="empty-catalog">✦ No docs indexed yet</div>', unsafe_allow_html=True)

    st.markdown('<div class="neon-divider"></div>', unsafe_allow_html=True)

    # ── Add document ───────────────────────────────────────────────────────
    st.markdown('<div class="sidebar-label">➕ Add to Knowledge Base</div>', unsafe_allow_html=True)
    add_tab = st.radio("Source", ["📎 Upload file", "✏️ Paste text"], horizontal=True, label_visibility="collapsed")

    if "Upload" in add_tab:
        uploaded = st.file_uploader(
            "Drop your file",
            type=["txt","md","pdf","docx","doc","xlsx","xls","csv","pptx","py","js","json","html"],
            label_visibility="collapsed"
        )
        title_uf = st.text_input("Title (optional)", key="title_uf", placeholder="Leave blank = filename")
        if st.button("⚡ Index File", key="add_file", use_container_width=True):
            if uploaded is None:
                st.error("Upload a file first!")
            else:
                with st.spinner("✨ Extracting & indexing…"):
                    try:
                        text = extract_text(uploaded)
                        if not text.strip():
                            st.error("No text found in file.")
                        else:
                            fname  = title_uf.strip() or uploaded.name
                            wcount = len(text.split())
                            added  = ingest_text(fname, text)
                            st.session_state.doc_meta[fname] = {"word_count": wcount, "chunks": added}
                            st.success(f"✓ {fname} — {added} chunks ready")
                            st.rerun()
                    except Exception as e:
                        st.error(f"Failed: {e}")
    else:
        paste_title = st.text_input("Title *", key="paste_title", placeholder="e.g. My notes")
        paste_text  = st.text_area("Content", key="paste_content", height=120,
                                   placeholder="Paste articles, docs, notes…")
        if st.button("⚡ Index Text", key="add_paste", use_container_width=True):
            if not paste_text.strip():
                st.error("Nothing to index!")
            elif not paste_title.strip():
                st.error("Give it a title!")
            else:
                with st.spinner("✨ Indexing…"):
                    wcount = len(paste_text.split())
                    added  = ingest_text(paste_title.strip(), paste_text)
                    st.session_state.doc_meta[paste_title.strip()] = {"word_count": wcount, "chunks": added}
                st.success(f"✓ {paste_title} — {added} chunks ready")
                st.rerun()

    st.markdown('<div class="neon-divider"></div>', unsafe_allow_html=True)

    # ── Memory ─────────────────────────────────────────────────────────────
    turns = len(st.session_state.chat_history)
    st.markdown(
        f'<div class="mem-badge"><span class="mem-dot"></span>'
        f'Memory &nbsp;·&nbsp; {turns} message{"s" if turns != 1 else ""}</div>',
        unsafe_allow_html=True
    )
    st.markdown("<div style='height:10px'></div>", unsafe_allow_html=True)

    ca, cb = st.columns(2)
    with ca:
        if st.button("🗑 Chat", use_container_width=True):
            st.session_state.chat_history = []
            st.rerun()
    with cb:
        if st.button("🗑 All", use_container_width=True, help="Clear chat + all docs"):
            for k in ["chat_history","rag_chunks","rag_tfidf","doc_list","doc_meta"]:
                st.session_state[k] = [] if k != "rag_idf" else {}
            st.session_state.rag_idf = {}
            st.rerun()

    st.markdown('<div class="neon-divider"></div>', unsafe_allow_html=True)

    # ── Settings ───────────────────────────────────────────────────────────
    with st.expander("⚙️ Settings"):
        st.session_state.top_k = st.slider(
            "Max total chunks (top-k)", 1, 20, st.session_state.top_k,
            help="Every doc gets at least 1 chunk guaranteed; remaining slots filled by score."
        )
        n_docs = len(st.session_state.doc_list)
        if n_docs > 0:
            st.markdown(
                f'<div style="font-family:monospace;font-size:10px;color:rgba(34,211,238,0.7);margin-top:-6px;margin-bottom:6px">'
                f'✦ {n_docs} doc(s) — each gets ≥1 chunk guaranteed</div>',
                unsafe_allow_html=True
            )
        st.session_state.show_sources = st.toggle("Show source citations", value=st.session_state.show_sources)
        st.session_state.show_chunks  = st.toggle("Debug: show chunks", value=st.session_state.show_chunks)

# ═══════════════════════════════════════════════════════════════════════════════
#  Main area
# ═══════════════════════════════════════════════════════════════════════════════

rag_active = len(st.session_state.doc_list) > 0

# ── Header ────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="masthead-wrap">
  <div class="masthead-title">🤖 AI Chatbot</div>
  <div class="masthead-sub">Powered by Gemini · RAG · Persistent Memory</div>
</div>
""", unsafe_allow_html=True)

# Status row
col_s1, col_s2, col_s3 = st.columns([1,1,4])
with col_s1:
    if rag_active:
        st.markdown(
            f'<div class="status-badge rag-on"><span class="status-dot"></span>'
            f'RAG ON · {len(st.session_state.doc_list)} docs</div>',
            unsafe_allow_html=True
        )
    else:
        st.markdown(
            '<div class="status-badge rag-off"><span class="status-dot"></span>RAG STANDBY</div>',
            unsafe_allow_html=True
        )
with col_s2:
    turns = len(st.session_state.chat_history)
    st.markdown(
        f'<div class="status-badge rag-off"><span class="status-dot"></span>'
        f'{turns} msg{"s" if turns!=1 else ""}</div>',
        unsafe_allow_html=True
    )

st.markdown('<div class="neon-divider"></div>', unsafe_allow_html=True)

# ── Chat history ──────────────────────────────────────────────────────────────
if not st.session_state.chat_history:
    st.markdown("""
    <div class="empty-state">
      <div class="big-emoji">✦</div>
      <h3>What's on your mind?</h3>
      <p>Ask me anything. Add docs on the left to chat with your files using RAG magic ✨</p>
    </div>
    """, unsafe_allow_html=True)
else:
    for msg in st.session_state.chat_history:
        if msg["role"] == "user":
            st.markdown(f'<div class="user-bubble">{msg["content"]}</div>', unsafe_allow_html=True)
        else:
            cls = "bot-bubble error" if msg.get("error") else "bot-bubble"
            content = msg["content"].replace("\n", "<br>")
            st.markdown(f'<div class="{cls}">{content}</div>', unsafe_allow_html=True)

            sources = msg.get("sources", [])
            if sources and st.session_state.show_sources:
                chips = "".join(f'<span class="chip">✦ {s}</span>' for s in sources)
                st.markdown(f'<div class="source-chips">{chips}</div>', unsafe_allow_html=True)

            if msg.get("chunks") and st.session_state.show_chunks:
                with st.expander("🔬 Retrieved chunks (grouped by document)"):
                    # Group chunks by document
                    from collections import defaultdict as _dd
                    by_doc = _dd(list)
                    for ch in msg["chunks"]:
                        by_doc[ch["doc"]].append(ch)
                    for doc_name, doc_chunks in by_doc.items():
                        st.markdown(
                            f'<div style="font-family:monospace;font-size:11px;color:#22d3ee;'
                            f'border-bottom:1px solid rgba(34,211,238,0.2);padding-bottom:4px;margin-bottom:8px">'
                            f'📄 {doc_name} — {len(doc_chunks)} chunk(s)</div>',
                            unsafe_allow_html=True
                        )
                        for i, ch in enumerate(doc_chunks, 1):
                            st.markdown(f"&nbsp;&nbsp;**Chunk {i}** — score: `{ch['score']:.3f}`")
                            st.caption("&nbsp;&nbsp;" + ch["text"][:350] + ("…" if len(ch["text"]) > 350 else ""))

# ── Chat input ────────────────────────────────────────────────────────────────
query = st.chat_input("Ask anything… ✨")

if query:
    if not st.session_state.api_key:
        st.error("⚠️ Add your Gemini API key in the sidebar first!")
    else:
        st.session_state.chat_history.append({"role": "user", "content": query})

        hits = _query_rag(query) if rag_active else []

        if hits:
            ctx = "\n".join(
                [f'=== DOCUMENT CONTEXT ==='] +
                [f'\n[Excerpt {i} from "{h["doc"]}"]\n{h["text"]}' for i, h in enumerate(hits, 1)] +
                ['\n=== END CONTEXT ===']
            )
            system = (
                ctx + "\n\nYou are a helpful AI assistant. Use the document context above to "
                "ground your answers. Cite sources like [1], [2] inline. Supplement with general "
                "knowledge if needed but don't misattribute it. Be clear and concise."
            )
        else:
            system = "You are a helpful, friendly AI assistant. Answer clearly and concisely."

        messages = [{"role": m["role"], "content": m["content"]} for m in st.session_state.chat_history]

        with st.spinner("✨ Thinking…"):
            reply = call_gemini(messages, system, st.session_state.api_key)

        st.session_state.chat_history.append({
            "role":    "assistant",
            "content": reply,
            "sources": list({h["doc"] for h in hits}),
            "chunks":  hits,
            "error":   reply.startswith("❌"),
        })
        st.rerun()
